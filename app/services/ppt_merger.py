from __future__ import annotations

import asyncio
from copy import deepcopy
from pathlib import Path

from pptx import Presentation

from app.utils.logger import get_logger

logger = get_logger(__name__)


def merge_ppt_files(files: list[str], *, output_path: str) -> str:
    """
    Merge multiple PPTX files into one deck.

    Note: python-pptx has no official slide-clone API; this uses a common XML element-copy approach.
    """
    if not files:
        raise ValueError("No PPT files to merge.")

    dest = Presentation(files[0])
    blank_layout = dest.slide_layouts[6] if len(dest.slide_layouts) > 6 else dest.slide_layouts[-1]

    for src_path in files[1:]:
        src = Presentation(src_path)
        for slide in src.slides:
            new_slide = dest.slides.add_slide(blank_layout)
            for shape in slide.shapes:
                el = shape.element
                new_el = deepcopy(el)
                new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")  # type: ignore[attr-defined]

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    dest.save(str(out))
    return str(out)


async def merge_ppt_files_async(files: list[str], *, output_path: str) -> str:
    return await asyncio.to_thread(merge_ppt_files, files, output_path=output_path)

