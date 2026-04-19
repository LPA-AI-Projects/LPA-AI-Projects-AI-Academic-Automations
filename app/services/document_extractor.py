from __future__ import annotations

import asyncio
from io import BytesIO
import re
from typing import Iterable

from app.utils.logger import get_logger

logger = get_logger(__name__)


def _join_nonempty(parts: Iterable[str]) -> str:
    cleaned = [p.strip() for p in parts if isinstance(p, str) and p.strip()]
    return "\n\n".join(cleaned).strip()


def extract_pdf_text(file_bytes: bytes) -> str:
    """
    Extract text from PDF bytes using PyMuPDF (fitz).
    Synchronous; call via asyncio.to_thread for non-blocking usage.
    """
    import fitz  # PyMuPDF

    doc = fitz.open(stream=file_bytes, filetype="pdf")
    try:
        pages: list[str] = []
        for page in doc:
            pages.append(page.get_text("text"))
        return _join_nonempty(pages)
    finally:
        doc.close()


def extract_ppt_text(file_bytes: bytes) -> str:
    """
    Extract text from PPTX bytes using python-pptx.
    Each slide is prefixed with ``--- Slide N ---`` so downstream code can split
    or score slides per module. Synchronous; call via asyncio.to_thread.
    """
    from io import BytesIO

    from pptx import Presentation

    prs = Presentation(BytesIO(file_bytes))
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_chunks: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_chunks.append(str(shape.text))
        body = _join_nonempty(slide_chunks)
        parts.append(f"--- Slide {idx} ---\n{body}" if body else f"--- Slide {idx} ---\n")
    return "\n\n".join(parts).strip()


_SLIDE_DELIM_RE = re.compile(r"(?m)^--- Slide (\d+) ---\s*$")

# Minimal English stopwords for keyword overlap scoring (module vs PPT slide text).
_PPT_SLICE_STOPWORDS = frozenset(
    {
        "the",
        "and",
        "for",
        "are",
        "but",
        "not",
        "you",
        "all",
        "can",
        "her",
        "was",
        "one",
        "our",
        "out",
        "day",
        "get",
        "has",
        "him",
        "his",
        "how",
        "its",
        "may",
        "new",
        "now",
        "old",
        "see",
        "two",
        "who",
        "way",
        "she",
        "use",
        "many",
        "then",
        "them",
        "these",
        "some",
        "what",
        "which",
        "when",
        "will",
        "with",
        "have",
        "this",
        "that",
        "from",
        "they",
        "been",
        "into",
        "more",
        "than",
        "also",
        "only",
        "such",
        "other",
        "about",
        "after",
        "module",
        "topic",
        "topics",
    }
)


def split_ppt_text_into_slide_blocks(ppt_text: str) -> list[tuple[int, str]]:
    """
    Parse text produced by extract_ppt_text into (slide_number, body) pairs in order.
    If no slide markers are present, returns a single block (1, full_text).
    """
    text = (ppt_text or "").strip()
    if not text:
        return []
    matches = list(_SLIDE_DELIM_RE.finditer(text))
    if not matches:
        return [(1, text)]
    blocks: list[tuple[int, str]] = []
    if matches[0].start() > 0:
        preamble = text[: matches[0].start()].strip()
    else:
        preamble = ""
    for i, m in enumerate(matches):
        num = int(m.group(1))
        start = m.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        body = text[start:end].strip()
        if i == 0 and preamble:
            body = (preamble + "\n\n" + body).strip()
        blocks.append((num, body))
    return blocks


def _keywords_for_module_slice(module_name: str, module_text: str, *, max_chars: int = 6000) -> set[str]:
    blob = f"{module_name}\n{(module_text or '')[:max_chars]}".lower()
    words = re.findall(r"[a-z][a-z\-]{2,}", blob)
    return {w for w in words if w not in _PPT_SLICE_STOPWORDS}


def slice_instructor_ppt_for_module(
    instructor_text: str | None,
    module_name: str,
    module_text: str,
    *,
    max_chars: int = 150_000,
) -> str | None:
    """
    Return PPT text most relevant to one outline module using keyword overlap per slide.
    Falls back to the full (truncated) deck when markers are missing or no slide scores.
    """
    raw = (instructor_text or "").strip()
    if not raw:
        return None
    blocks = split_ppt_text_into_slide_blocks(raw)
    if len(blocks) <= 1 and "--- Slide " not in raw:
        return raw[:max_chars]

    keys = _keywords_for_module_slice(module_name, module_text)
    if not keys:
        return raw[:max_chars]

    scores: list[tuple[int, int, str]] = []
    for num, body in blocks:
        low = body.lower()
        score = sum(low.count(k) for k in keys)
        scores.append((score, num, body))

    max_score = max((s[0] for s in scores), default=0)
    if max_score == 0:
        return raw[:max_chars]

    # Include neighboring slides around any slide with a positive score.
    want: set[int] = set()
    for score, num, _ in scores:
        if score <= 0:
            continue
        want.update({num - 1, num, num + 1})
    by_num = {num: body for num, body in blocks}
    ordered_nums = sorted(n for n in want if n in by_num)
    if not ordered_nums:
        return raw[:max_chars]

    out_parts: list[str] = []
    total = 0
    for n in ordered_nums:
        chunk = f"--- Slide {n} ---\n{by_num[n]}".strip()
        if total + len(chunk) + 2 > max_chars:
            break
        out_parts.append(chunk)
        total += len(chunk) + 2
    joined = "\n\n".join(out_parts).strip()
    return joined if joined else raw[:max_chars]


async def extract_pdf_text_async(file_bytes: bytes) -> str:
    return await asyncio.to_thread(extract_pdf_text, file_bytes)


async def extract_ppt_text_async(file_bytes: bytes) -> str:
    return await asyncio.to_thread(extract_ppt_text, file_bytes)


def extract_pdf_module_rows(file_bytes: bytes) -> list[dict[str, str]]:
    """
    Extract module rows from table-like course outlines (Sno | Modules | Topics | Exercises).
    Returns [] when no usable table rows are found.
    """
    try:
        import pdfplumber
    except Exception:
        return []

    def _norm(v: object) -> str:
        return str(v or "").strip()

    def _is_header_row(row: list[str]) -> bool:
        joined = " | ".join(row).lower()
        return (
            "sno" in joined
            and "module" in joined
            and ("topic" in joined or "content" in joined)
            and ("exercise" in joined or "activity" in joined)
        )

    modules: list[dict[str, str]] = []
    with pdfplumber.open(BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            for table in page.extract_tables() or []:
                if not table:
                    continue
                rows: list[list[str]] = [[_norm(c) for c in (r or [])] for r in table]
                if not rows:
                    continue

                start_idx = -1
                for i, row in enumerate(rows):
                    if _is_header_row(row):
                        start_idx = i + 1
                        break
                # Only parse table rows when a valid module-table header exists.
                if start_idx < 0:
                    continue

                for row in rows[start_idx:]:
                    if len(row) < 3:
                        continue
                    sno = _norm(row[0])
                    if not sno:
                        continue
                    m = re.match(r"^0?(\d{1,2})$", sno)
                    if not m:
                        continue
                    module_num = int(m.group(1))
                    module_name = _norm(row[1])
                    topics = _norm(row[2]) if len(row) > 2 else ""
                    exercises = _norm(row[3]) if len(row) > 3 else ""
                    if not module_name:
                        continue
                    if module_name.strip().lower() in {"module", "modules"}:
                        continue
                    module_text_parts = [
                        f"Module {module_num}: {module_name}",
                        f"Topics:\n{topics}" if topics else "",
                        f"Exercises:\n{exercises}" if exercises else "",
                    ]
                    module_text = "\n\n".join([p for p in module_text_parts if p]).strip()
                    modules.append(
                        {
                            "module_name": f"Module {module_num}: {module_name}",
                            "module_text": module_text,
                        }
                    )

    # Deduplicate by module number, keep first occurrence.
    deduped: list[dict[str, str]] = []
    seen: set[int] = set()
    for item in modules:
        n_match = re.search(r"Module\s+(\d+)", item.get("module_name", ""))
        if not n_match:
            continue
        n = int(n_match.group(1))
        if n in seen:
            continue
        seen.add(n)
        deduped.append(item)
    return deduped


async def extract_pdf_module_rows_async(file_bytes: bytes) -> list[dict[str, str]]:
    return await asyncio.to_thread(extract_pdf_module_rows, file_bytes)

